from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from burger_types import BURGER_TYPES
import os

def create_burger_presentation():
    """Create a PowerPoint presentation for the Food-Graph project"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(255, 87, 34)  # Orange
    SUBTITLE_COLOR = RGBColor(66, 66, 66)  # Dark gray
    
    print("Creating PowerPoint presentation...")
    print("=" * 80)
    
    # Slide 1: Title Slide
    print("\n✓ Creating Slide 1: Title Slide")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "🍔 Food-Graph Project"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Graph-Based Burger Calorie Analysis System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = SUBTITLE_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add author info
    author_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    author_frame = author_box.text_frame
    author_frame.text = "Data Structures & Graph Theory Application"
    author_para = author_frame.paragraphs[0]
    author_para.font.size = Pt(18)
    author_para.font.italic = True
    author_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    print("✓ Creating Slide 2: Project Overview")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A nutrition analysis system that uses graph data structures to:"
    
    points = [
        "Model relationships between burgers and ingredients",
        "Calculate total calories through graph traversal",
        "Enable complex queries on nutritional data",
        "Visualize calorie flow and ingredient relationships"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 3: Key Features
    print("✓ Creating Slide 3: Key Features")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    features = [
        ("8 Different Burger Types", "With unique ingredient combinations"),
        ("30+ Ingredients", "Categorized by type (buns, proteins, cheese, etc.)"),
        ("Graph-Based Architecture", "Using NetworkX for relationship modeling"),
        ("Advanced Queries", "Find burgers by calorie range, ingredients, etc."),
        ("Visual Representations", "Multiple chart types and flow diagrams")
    ]
    
    for idx, (feature, desc) in enumerate(features):
        if idx == 0:
            tf.text = f"{feature}"
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = feature
        
        p.font.size = Pt(22)
        p.font.bold = True
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.level = 1
        p_desc.font.size = Pt(18)
    
    # Slide 4: Graph Structure
    print("✓ Creating Slide 4: Graph Structure")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Graph Data Structure"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Directed Graph Implementation:"
    
    structure_points = [
        "Nodes: Burgers and Ingredients (35 total)",
        "Edges: Burger → Ingredient relationships (59 connections)",
        "Edge Weights: Calorie contribution values",
        "Graph Statistics: 8 burgers, 27 unique ingredients",
        "Average: 7.4 ingredients per burger"
    ]
    
    for point in structure_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Add graph network image if available
    graph_img_path = "burger_images/networks/burger_graph_network.png"
    if os.path.exists(graph_img_path):
        slide.shapes.add_picture(graph_img_path, Inches(5.5), Inches(2), height=Inches(4))
    
    # Slide 5: Burger Types & Calories
    print("✓ Creating Slide 5: Burger Types & Calories")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Burger Types Ranked by Calories"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Add comparison chart
    comparison_img_path = "burger_images/comparisons/all_burgers_comparison.png"
    if os.path.exists(comparison_img_path):
        slide.shapes.add_picture(comparison_img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    
    # Slide 6: Calorie Analysis
    print("✓ Creating Slide 6: Calorie Analysis")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Calorie Analysis Results"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Calculate stats
    burgers_data = [(burger.name, burger.calculate_calories()) 
                   for burger in BURGER_TYPES.values()]
    burgers_data.sort(key=lambda x: x[1])
    
    lowest = burgers_data[0]
    highest = burgers_data[-1]
    avg_calories = sum(cal for _, cal in burgers_data) / len(burgers_data)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Lowest Calorie: {lowest[0]} - {lowest[1]} cal"
    tf.paragraphs[0].font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = f"Highest Calorie: {highest[0]} - {highest[1]} cal"
    p.font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = f"Average Calories: {avg_calories:.0f} cal"
    p.font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = f"Calorie Range: {highest[1] - lowest[1]} cal difference"
    p.font.size = Pt(22)
    
    # Add distribution chart
    dist_img_path = "burger_images/distributions/calorie_distribution.png"
    if os.path.exists(dist_img_path):
        slide.shapes.add_picture(dist_img_path, Inches(5), Inches(2.5), height=Inches(4))
    
    # Slide 7: Calorie Flow Diagram
    print("✓ Creating Slide 7: Calorie Flow Diagram")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Calorie Flow Visualization"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Add flow diagram if available
    flow_img_path = "burger_images/flow_diagrams/double_deluxe_calorie_flow.png"
    if os.path.exists(flow_img_path):
        slide.shapes.add_picture(flow_img_path, Inches(0.5), Inches(1.5), width=Inches(9))
    
    # Slide 8: Individual Burger Example
    print("✓ Creating Slide 8: Individual Burger Example")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Classic Burger - Layer Breakdown"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    # Add burger diagram
    burger_img_path = "burger_images/individual_diagrams/classic_diagram.png"
    if os.path.exists(burger_img_path):
        slide.shapes.add_picture(burger_img_path, Inches(1), Inches(1.5), height=Inches(5.5))
    
    # Slide 9: Technical Implementation
    print("✓ Creating Slide 9: Technical Implementation")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Implementation"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technologies Used:"
    
    tech_points = [
        "Python 3.12 - Core programming language",
        "NetworkX - Graph data structure library",
        "Matplotlib - Data visualization framework",
        "Object-Oriented Design - Clean class hierarchies",
        "Modular Architecture - Separated concerns"
    ]
    
    for point in tech_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 10: Graph Operations
    print("✓ Creating Slide 10: Graph Operations")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Graph Operations & Queries"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Supported Operations:"
    
    ops_points = [
        "Calculate total calories by graph traversal",
        "Find burgers containing specific ingredients",
        "Query burgers within calorie ranges",
        "Compare two burgers (ingredients & calories)",
        "Get graph statistics and insights"
    ]
    
    for point in ops_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(20)
    
    # Slide 11: Key Concepts Demonstrated
    print("✓ Creating Slide 11: Key Concepts")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Concepts Demonstrated"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    concepts = [
        ("Graph Data Structures", "Modeling complex relationships"),
        ("Object-Oriented Design", "Clean class hierarchies"),
        ("Data Aggregation", "Calculating totals via traversal"),
        ("Data Visualization", "Multiple chart types"),
        ("Nutritional Analysis", "Comparing food items")
    ]
    
    for idx, (concept, desc) in enumerate(concepts):
        if idx == 0:
            tf.text = f"{concept}"
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = concept
        
        p.font.size = Pt(22)
        p.font.bold = True
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.level = 1
        p_desc.font.size = Pt(18)
    
    # Slide 12: Project Structure
    print("✓ Creating Slide 12: Project Structure")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Structure"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    structure = [
        "ingredients.py - Ingredient definitions",
        "burger_types.py - Burger type definitions",
        "burger_graph.py - Graph implementation",
        "calorie_calculator.py - Analysis engine",
        "visualize_burgers.py - Visualization generator",
        "burger_flow_diagram.py - Flow diagram creator"
    ]
    
    tf.text = structure[0]
    for point in structure[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
    
    # Slide 13: Future Enhancements
    print("✓ Creating Slide 13: Future Enhancements")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    title.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Potential Extensions:"
    
    future_points = [
        "Add more nutritional data (protein, fat, carbs, sodium)",
        "Implement dietary restriction filters (vegan, gluten-free)",
        "Create interactive web interface",
        "Add recommendation system based on preferences",
        "Include pricing and cost optimization",
        "Expand to other food categories"
    ]
    
    for point in future_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 14: Conclusion
    print("✓ Creating Slide 14: Conclusion")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add conclusion text
    conclusion_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
    conclusion_frame = conclusion_box.text_frame
    conclusion_frame.word_wrap = True
    
    conclusion_text = [
        "Successfully demonstrated practical application of graph data structures",
        "Created comprehensive visualization and analysis system",
        "Showcased object-oriented design and modular architecture",
        "Built scalable solution for nutritional analysis"
    ]
    
    for idx, text in enumerate(conclusion_text):
        if idx == 0:
            conclusion_frame.text = f"• {text}"
            p = conclusion_frame.paragraphs[0]
        else:
            p = conclusion_frame.add_paragraph()
            p.text = f"• {text}"
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.LEFT
    
    # Slide 15: Thank You
    print("✓ Creating Slide 15: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add thank you message
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(60)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = TITLE_COLOR
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Questions & Discussion"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.italic = True
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "Food-Graph_Presentation.pptx"
    prs.save(output_file)
    
    print("\n" + "=" * 80)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: 15")
    print("=" * 80)
    print("\nPresentation includes:")
    print("  • Title and overview slides")
    print("  • Technical implementation details")
    print("  • Visual representations and charts")
    print("  • Graph structure explanation")
    print("  • Calorie analysis results")
    print("  • Future enhancements")
    print("\nReady to present to your supervisor!")

if __name__ == "__main__":
    create_burger_presentation()
